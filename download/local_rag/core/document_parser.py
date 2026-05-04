"""
文档解析器 — 支持多种格式
================================
借鉴 WeKnora 的 docparser 设计，支持 PDF/DOCX/TXT/MD/HTML 等
常见文档格式的解析，输出统一的文本内容。
"""

import os
from pathlib import Path
from typing import Optional

from loguru import logger


class DocumentParser:
    """多格式文档解析器"""

    SUPPORTED_EXTENSIONS = {
        ".pdf", ".docx", ".doc", ".txt", ".md", ".html", ".htm",
        ".pptx", ".csv", ".xlsx", ".xls", ".json",
    }

    @staticmethod
    def is_supported(file_path: str) -> bool:
        """检查文件格式是否支持"""
        ext = Path(file_path).suffix.lower()
        return ext in DocumentParser.SUPPORTED_EXTENSIONS

    @staticmethod
    async def parse(file_path: str) -> dict:
        """
        解析文档，返回统一格式的解析结果

        Args:
            file_path: 文件路径

        Returns:
            {
                "title": str,
                "content": str,
                "file_type": str,
                "file_name": str,
                "file_size": int,
                "metadata": dict,
            }
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")

        ext = path.suffix.lower()
        file_name = path.name
        file_size = path.stat().st_size

        parsers = {
            ".pdf": DocumentParser._parse_pdf,
            ".docx": DocumentParser._parse_docx,
            ".doc": DocumentParser._parse_docx,
            ".txt": DocumentParser._parse_text,
            ".md": DocumentParser._parse_text,
            ".html": DocumentParser._parse_html,
            ".htm": DocumentParser._parse_html,
            ".pptx": DocumentParser._parse_pptx,
            ".csv": DocumentParser._parse_csv,
            ".xlsx": DocumentParser._parse_xlsx,
            ".xls": DocumentParser._parse_xlsx,
            ".json": DocumentParser._parse_json,
        }

        parser = parsers.get(ext)
        if not parser:
            raise ValueError(f"不支持的文件格式: {ext}")

        content = await parser(file_path)

        # 从内容推断标题
        title = DocumentParser._extract_title(content, file_name)

        return {
            "title": title,
            "content": content,
            "file_type": ext.lstrip("."),
            "file_name": file_name,
            "file_size": file_size,
            "metadata": {},
        }

    @staticmethod
    async def _parse_pdf(file_path: str) -> str:
        """解析 PDF 文件"""
        import pypdf

        texts = []
        with open(file_path, "rb") as f:
            reader = pypdf.PdfReader(f)
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    texts.append(text.strip())
        return "\n\n".join(texts)

    @staticmethod
    async def _parse_docx(file_path: str) -> str:
        """解析 DOCX 文件"""
        from docx import Document

        doc = Document(file_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)

    @staticmethod
    async def _parse_text(file_path: str) -> str:
        """解析纯文本/Markdown 文件"""
        import aiofiles

        async with aiofiles.open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return await f.read()

    @staticmethod
    async def _parse_html(file_path: str) -> str:
        """解析 HTML 文件"""
        from bs4 import BeautifulSoup

        import aiofiles
        async with aiofiles.open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            html = await f.read()
        soup = BeautifulSoup(html, "html.parser")

        # 移除脚本和样式
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        return soup.get_text(separator="\n", strip=True)

    @staticmethod
    async def _parse_pptx(file_path: str) -> str:
        """解析 PPTX 文件"""
        from pptx import Presentation

        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                texts.append("\n".join(slide_texts))
        return "\n\n---\n\n".join(texts)

    @staticmethod
    async def _parse_csv(file_path: str) -> str:
        """解析 CSV 文件"""
        import csv

        rows = []
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(" | ".join(row))
        return "\n".join(rows)

    @staticmethod
    async def _parse_xlsx(file_path: str) -> str:
        """解析 XLSX/XLS 文件"""
        from openpyxl import load_workbook

        wb = load_workbook(file_path, read_only=True)
        texts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    texts.append(" | ".join(cells))
        wb.close()
        return "\n".join(texts)

    @staticmethod
    async def _parse_json(file_path: str) -> str:
        """解析 JSON 文件"""
        import json
        import aiofiles

        async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
            content = await f.read()
        data = json.loads(content)
        return json.dumps(data, ensure_ascii=False, indent=2)

    @staticmethod
    def _extract_title(content: str, file_name: str) -> str:
        """从内容中提取标题"""
        for line in content.split("\n")[:5]:
            line = line.strip()
            if line.startswith("# "):
                return line[2:].strip()
            if line and len(line) < 100:
                return line
        return Path(file_name).stem
